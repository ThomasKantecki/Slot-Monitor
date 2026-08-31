#!/usr/bin/env python3
"""The story deck: how the provider map came to be, and why its numbers hold.

One content model, two renderers:
  ~/Desktop/Provider-Map-Deck.pptx  - the deck to present (16:9, editable)
  ~/Desktop/Provider-Map-Deck.html  - pixel-true twin for rehearsal (this Mac
                                      has no PowerPoint; the twin embeds the
                                      map's actual webfonts)

Every number that exists in the map's data artifacts is read from them here,
not typed. Screenshots are optional: pass paths via MAP_COUNTY_SHOT /
MAP_ZIP_SHOT env vars; slides degrade gracefully without them.

PPTX fonts are portability-safe (Courier New / Helvetica Neue) because
python-pptx cannot embed fonts and the presenting machine is unknown.
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DESK = Path.home() / "Desktop"

# ---- the map's design tokens (src/render.js :root) --------------------------
CREAM = "F5F1E8"
CREAM90 = "EEE9DC"
NAVY = "14233E"
MUTE = "41506C"
FAINT = "63748C"
AH = "005C99"
OH = "B20838"
WHITE = "FFFFFF"
BLACK = "000000"
BAND_META = "9DC3DF"

MONO = "Courier New"
DISP = "Helvetica Neue"

AH_LOGO = ROOT / "assets" / "adventhealth-logo.png"
OH_LOGO = ROOT / "assets" / "orlandohealth-logo.png"
COUNTY_SHOT = os.environ.get("MAP_COUNTY_SHOT", "")
ZIP_SHOT = os.environ.get("MAP_ZIP_SHOT", "")

# ---- live numbers from the artifacts ----------------------------------------
byzip = json.loads((ROOT / "data" / "providers-by-zip.json").read_text())
bycounty = json.loads((ROOT / "data" / "providers-by-county.json").read_text())
TOT_AH, TOT_OH = byzip["totals"]["ah"], byzip["totals"]["oh"]
N_ZIPS = len(byzip["zips"])
N_COUNTIES = len(bycounty["zips"])
N_SPECS = len(byzip["specialties"])
county_rows = sorted(bycounty["zips"].items(), key=lambda kv: -(kv[1]["ah"] + kv[1]["oh"]))
head2head = [s for s in byzip["specialties"] if s["ah"] > 0 and s["oh"] > 0]
head2head = sorted(head2head, key=lambda s: -(s["ah"] + s["oh"]))[:7]
fmt = lambda n: f"{n:,}"

# =============================================================================
# content model
# =============================================================================
# Each slide: {"band": str, "meta": str, "layout": ..., ...layout-specific keys}

SLIDES = [
    {
        "layout": "title",
        "mark": ("PROVIDER ", "MAP"),
        "subtitle": "Bookable provider coverage by area: AdventHealth vs Orlando Health",
        "date": "August 31, 2026  ·  prepared by Thomas Kantecki",
    },
    {
        "band": "WHERE WE STARTED: APPOINTMENT SLOTS",
        "layout": "stats+bullets",
        "stats": [
            ("AH", "64", "specialties bookable online", "41,541 slots captured in one sweep"),
            ("OH", "9", "specialties bookable online", "892 slots because most specialties are phone only"),
        ],
        "bullets": [
            "We scraped anonymous MyChart open scheduling for both systems and captured all the slots they publish.",
            "The problem is that most specialties never publish slots online. You have to call the office to get an appointment, so those specialties show up as zero.",
            "That means slot counts mostly tell you how each website is set up. They say very little about how many providers a system actually has.",
        ],
    },
    {
        "band": "WHY SLOT COUNTS FAILED AS A METRIC",
        "layout": "bigstat+bullets",
        "big": ("~85%", "of each system's providers never publish a single online slot"),
        "bullets": [
            "Only 15 to 17 percent of providers on either side expose online booking (OH 17%, AH 15%).",
            "Orlando Health gates most online flows behind screening questions. AdventHealth exposes about 7x more specialties online. Slot totals reflect those choices.",
            "A metric that can't see 85% of the market can't support real decisions.",
        ],
        "verdict": "VERDICT: slot counts measure online scheduling policy. We needed a capacity metric.",
    },
    {
        "band": "THE PIVOT: PROVIDERS PER SPECIALTY PER ZIP",
        "layout": "cards",
        "lead": "Count the clinicians each system employs in bookable clinic specialties, organized by where they practice.",
        "cards": [
            ("COMPLETE", "Every employed clinician appears in their system's own directory, so nothing is hidden by website settings."),
            ("COMPARABLE", "Both sides get the same definition: employed, in a bookable specialty, counted at one practice location."),
            ("VERIFIABLE", "Every provider has a federal NPI number that can be checked against the national registry at any time."),
            ("GEOGRAPHIC", "Rolls up from ZIP code to county to statewide and drills down to individual named providers."),
        ],
    },
    {
        "band": "WHAT WE TRIED FIRST",
        "layout": "sequence",
        "steps": [
            ("BLOCKED", "orlandohealth.com returns HTTP 429 on every page. The whole site sits behind a bot wall, including robots.txt."),
            ("THE DETOUR", "We rebuilt from federal CMS data instead. A CMS enrollment record turned out to be a billing permissions list that says almost nothing about where someone works. Placement accuracy came out to 42.7% and one clinician showed up in 17 different ZIP codes."),
            ("SCRAPPED", "After a full day of pipeline work, name matching and manual review lists, we deleted all of it."),
        ],
        "verdict": "LESSON: piecing together public data from federal sources and scattered databases doesn't work. The data needs to come from each system's own directory.",
    },
    {
        "band": "BREAKTHROUGH 1: ORLANDO HEALTH",
        "layout": "stats+bullets",
        "stats": [
            ("OH", "8", "requests", "for their complete directory"),
            ("OH", "7,732", "provider records", "with NPIs, employment flags and geocoded locations"),
        ],
        "bullets": [
            "Their find a doctor page is an empty shell. The actual doctor list comes from a hosted search service that answers queries from any visitor's browser.",
            "We queried that service directly using the same request the page makes with a bigger page size.",
            "The records contain more than the site ever displays, including the NPI, an employment flag and every practice location with coordinates.",
        ],
    },
    {
        "band": "BREAKTHROUGH 2: ADVENTHEALTH",
        "layout": "twopanel",
        "left": ("HOW WE GOT THE DATA", [
            "Our site has no separate data service. Every page is rendered on the server behind Akamai bot protection, so outside scripts get blocked.",
            "A real Chrome session passes the bot wall the normal way, so we ran the scrape inside the browser itself.",
            "It walked all 798 result pages with zero failures and captured the name, specialty, practice group and address for all 9,551 providers in the network.",
        ]),
        "right": ("HOW WE KNOW WE GOT EVERYONE", [
            "The public sitemap lists every doctor page for search engines. That is 10,160 URLs and each one ends in the doctor's NPI number.",
            "That gave us an independent roster to check the scrape against.",
            "All 2,417 employed providers on the map also appear in the sitemap, and the scrape covered 99.75% of the live listing.",
        ]),
    },
    {
        "band": "ONE METHODOLOGY, BOTH SIDES",
        "layout": "cards+proof",
        "cards": [
            ("EMPLOYED ONLY", "Orlando Health has an employment flag in their data. For AdventHealth we use AdventHealth branded practice groups. Independent practices are excluded on both sides."),
            ("BOOKABLE ONLY", "Hospital staff and support roles are excluded from both sides. Both systems have these people but only Orlando Health publishes them as named providers, so counting them would create false zeros for AdventHealth."),
            ("ONE VOCABULARY", "The directories name specialties differently. OBGYN and Obstetrics and Gynecology are the same thing, so the naming is merged into one list."),
            ("ONE LOCATION", "Each clinician counts once at their primary practice site. 534 multi site Orlando Health clinicians were reduced to one location."),
        ],
        "proof": "PROOF THE FILTERS MATTERED: the bookable filter moved Orlando Health from 3,995 down to 2,728 while AdventHealth stayed at 2,417. The inflation was entirely on their side.",
    },
    {
        "band": "WE ALSO FIXED THE SOURCES' OWN ERRORS",
        "layout": "cards",
        "lead": "Both directories publish mistakes. Every fix below is verified and documented in code.",
        "cards": [
            ("WRONG NPIs", "Four published NPIs were malformed or belonged to a different clinician. Correcting them against the federal registry also exposed three duplicate records, which we merged."),
            ("ZIP TYPOS", "29 locations had bad ZIP codes. When eight colleagues at the same clinic say 32765 and one says 32755, the odd one out is a typo, so we fixed it."),
            ("PHANTOM SITES", "30 anesthesiologists were listed at a staffing firm's corporate office in Fort Lauderdale. Nobody sees patients there, so those addresses no longer place anyone on the map."),
            ("DUPLICATES", "Some people were listed twice under a nickname or a stale record. Those merge by NPI with a same name guard."),
        ],
    },
    {
        "band": "THE RESULT",
        "layout": "result",
        "counties": county_rows[:8],
        "shot": COUNTY_SHOT,
    },
    {
        "band": "WHAT THE MAP ANSWERS",
        "layout": "h2h",
        "rows": head2head,
        "shot": ZIP_SHOT,
        "bullets": [
            "Shows where either system is thin in any specialty, down to the ZIP code.",
            "Shows which contested counties are close and which service lines drive the gap.",
            "Every number drills to a named provider at a real address.",
        ],
    },
    {
        "band": "APPENDIX: METHODOLOGY FINE PRINT",
        "layout": "threecol",
        "cols": [
            ("WHY PEOPLE ARE EXCLUDED", [
                "Exclusions reflect what each directory publishes, not what each system offers.",
                "Example: AdventHealth runs 40+ physical therapy locations but lists no individual PTs online. Counting PTs would show AdventHealth at zero, which is false.",
                "Excluded from both sides: hospital based specialties, radiology, pathology, PT, speech, nutrition, CRNAs and records with no stated specialty.",
            ]),
            ("DISCLOSED ASYMMETRIES", [
                "AdventHealth lists one site per clinician while Orlando Health lists all of them, so we count the primary site for both",
                "AdventHealth employment is inferred from practice group branding while Orlando Health's is an explicit flag",
                "The AdventHealth capture is a snapshot while Orlando Health can be re-pulled on demand",
            ]),
            ("EDGE CASES", [
                "9 clinicians appear in both directories and stay in both, matching what each system publishes",
                "The AdventHealth listing was captured twice in full and both passes returned the identical 9,551 providers",
                "Border ZIP codes are assigned to their majority county",
            ]),
        ],
    },
    {
        "layout": "title",
        "mark": ("THE ", "END"),
        "subtitle": "",
        "date": "",
    },
]

# =============================================================================
# PPTX renderer
# =============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
rgb = lambda h: RGBColor.from_string(h)


def box(slide, x, y, w, h, fill=None, line=None, line_w=2.0):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line)
        sh.line.width = Pt(line_w)
    return sh


def text(slide, x, y, w, h, runs, size=14, font=DISP, color=NAVY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.15, wrap=True):
    """runs: str, or list of paragraphs, each str or list of (txt, overrides)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        if i > 0:
            p.space_before = Pt(size * 0.45)
        for piece in (para if isinstance(para, list) else [(para, {})]):
            t, ov = piece if isinstance(piece, tuple) else (piece, {})
            r = p.add_run()
            r.text = t
            r.font.name = ov.get("font", font)
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.color.rgb = rgb(ov.get("color", color))
    return tb


def base(slide):
    box(slide, -0.06, -0.06, 13.45, 7.62, fill=CREAM)


def band(slide, title, meta=""):
    box(slide, 0.5, 0.42, 12.33, 0.62, fill=NAVY)
    text(slide, 0.78, 0.56, 9.6, 0.4, title, size=17, font=MONO, color=CREAM, bold=True)
    if meta:
        text(slide, 8.2, 0.6, 4.35, 0.36, meta, size=12.5, font=MONO, color=BAND_META,
             bold=True, align=PP_ALIGN.RIGHT)


def kicker(slide, x, y, w, t, color=FAINT, size=10.5):
    text(slide, x, y, w, 0.3, t.upper(), size=size, font=MONO, color=color, bold=True)


def logo(slide, which, x, y, h=0.42):
    p = AH_LOGO if which == "AH" else OH_LOGO
    ratios = {"AH": 462 / 112, "OH": 343 / 120}
    slide.shapes.add_picture(str(p), Inches(x), Inches(y), height=Inches(h),
                             width=Inches(h * ratios[which]))


def bullets_block(slide, x, y, w, items, size=19):
    paras = [[("-  ", {"font": MONO, "bold": True, "color": FAINT}), (b, {})] for b in items]
    text(slide, x, y, w, 4.5, paras, size=size, color=MUTE, leading=1.28)


for spec in SLIDES:
    slide = prs.slides.add_slide(BLANK)
    base(slide)
    layout = spec["layout"]

    if layout == "title":
        box(slide, 0, 0, 13.34, 1.9, fill=NAVY)
        mk = box(slide, 3.57, 2.5, 6.2, 1.1, fill=WHITE, line=BLACK, line_w=3.0)
        text(slide, 3.57, 2.76, 6.2, 0.6,
             [[(spec["mark"][0], {"color": NAVY}), (spec["mark"][1], {"color": AH})]],
             size=36, font=MONO, bold=True, align=PP_ALIGN.CENTER)
        logo(slide, "AH", 4.45, 4.1, h=0.58)
        text(slide, 6.6, 4.18, 0.3, 0.45, "×", size=24, font=MONO, color=FAINT, align=PP_ALIGN.CENTER)
        logo(slide, "OH", 7.05, 4.06, h=0.64)
        text(slide, 1.5, 5.25, 10.33, 0.55, spec["subtitle"], size=20, color=NAVY,
             align=PP_ALIGN.CENTER)
        text(slide, 1.5, 6.65, 10.33, 0.4, spec["date"], size=12.5, font=MONO, color=FAINT,
             align=PP_ALIGN.CENTER)
        continue

    band(slide, spec["band"], spec.get("meta", ""))

    if layout == "stats+bullets":
        y = 1.7
        for sys, num, cap, sub in spec["stats"]:
            box(slide, 0.5, y, 5.75, 2.3, fill=WHITE, line=BLACK)
            logo(slide, sys, 0.85, y + 0.3, h=0.48)
            text(slide, 3.0, y + 0.22, 3.0, 1.05, num, size=62, font=MONO,
                 color=(AH if sys == "AH" else OH), bold=True, align=PP_ALIGN.RIGHT)
            text(slide, 0.85, y + 1.18, 5.15, 0.38, cap.upper(), size=14, font=MONO,
                 color=FAINT, bold=True)
            text(slide, 0.85, y + 1.62, 5.15, 0.5, sub, size=15.5, color=MUTE)
            y += 2.62
        bullets_block(slide, 6.7, 1.78, 6.05, spec["bullets"], size=19)

    elif layout == "bigstat+bullets":
        box(slide, 0.5, 1.55, 5.75, 4.3, fill=WHITE, line=BLACK)
        text(slide, 0.7, 2.2, 5.35, 1.9, spec["big"][0], size=120, font=MONO, color=OH,
             bold=True, align=PP_ALIGN.CENTER)
        text(slide, 0.95, 4.4, 4.9, 1.2, spec["big"][1].upper(), size=15.5, font=MONO,
             color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        bullets_block(slide, 6.7, 1.65, 6.05, spec["bullets"], size=19)
        box(slide, 0.5, 6.25, 12.33, 0.9, fill=CREAM90, line=BLACK)
        text(slide, 0.8, 6.51, 11.7, 0.48, spec["verdict"], size=17, font=MONO,
             color=NAVY, bold=True)

    elif layout == "cards":
        if spec.get("lead"):
            text(slide, 0.5, 1.3, 12.3, 0.55, spec["lead"], size=21, color=NAVY, bold=True)
        y0 = 2.05 if spec.get("lead") else 1.6
        ch = 2.45 if spec.get("lead") else 2.65
        cw, gap = 6.02, 0.3
        for i, (t, b) in enumerate(spec["cards"]):
            x = 0.5 + (i % 2) * (cw + gap)
            y = y0 + (i // 2) * (ch + gap)
            box(slide, x, y, cw, ch, fill=WHITE, line=BLACK)
            box(slide, x, y, 0.18, ch, fill=(AH if i % 2 == 0 else OH))
            kicker(slide, x + 0.44, y + 0.26, cw - 0.75, t, color=NAVY, size=15.5)
            text(slide, x + 0.44, y + 0.74, cw - 0.8, ch - 0.9, b, size=16.5, color=MUTE)

    elif layout == "sequence":
        xs = [0.5, 4.7, 8.9]
        for i, (t, b) in enumerate(spec["steps"]):
            box(slide, xs[i], 1.5, 3.95, 4.55, fill=WHITE, line=BLACK)
            box(slide, xs[i], 1.5, 3.95, 0.64, fill=NAVY)
            text(slide, xs[i] + 0.24, 1.66, 3.6, 0.4, f"{i+1}  {t}", size=16.5, font=MONO,
                 color=CREAM, bold=True)
            text(slide, xs[i] + 0.3, 2.45, 3.4, 3.4, b, size=16.5, color=MUTE)
        box(slide, 0.5, 6.3, 12.33, 0.9, fill=CREAM90, line=BLACK)
        text(slide, 0.8, 6.56, 11.7, 0.48, spec["verdict"], size=17, font=MONO,
             color=NAVY, bold=True)

    elif layout == "twopanel":
        for i, (title, items) in enumerate([spec["left"], spec["right"]]):
            x = 0.5 + i * 6.3
            box(slide, x, 1.45, 6.02, 5.6, fill=WHITE, line=BLACK)
            box(slide, x, 1.45, 6.02, 0.64, fill=(AH if i else NAVY))
            text(slide, x + 0.3, 1.61, 5.5, 0.4, title, size=16.5, font=MONO, color=CREAM, bold=True)
            paras = [[("-  ", {"font": MONO, "bold": True, "color": FAINT}), (b, {})] for b in items]
            text(slide, x + 0.36, 2.5, 5.35, 4.4, paras, size=17, color=MUTE, leading=1.32)

    elif layout == "cards+proof":
        cw, ch, gap = 6.02, 2.15, 0.28
        for i, (t, b) in enumerate(spec["cards"]):
            x = 0.5 + (i % 2) * (cw + gap)
            y = 1.4 + (i // 2) * (ch + gap)
            box(slide, x, y, cw, ch, fill=WHITE, line=BLACK)
            box(slide, x, y, 0.18, ch, fill=(AH if i % 2 == 0 else OH))
            kicker(slide, x + 0.44, y + 0.22, cw - 0.75, t, color=NAVY, size=15)
            text(slide, x + 0.44, y + 0.66, cw - 0.8, ch - 0.8, b, size=15, color=MUTE)
        box(slide, 0.5, 6.05, 12.33, 1.15, fill=NAVY)
        text(slide, 0.8, 6.25, 11.73, 0.78, spec["proof"], size=16, font=MONO,
             color=CREAM, bold=True)

    elif layout == "threecol":
        for i, (t, items) in enumerate(spec["cols"]):
            x = 0.5 + i * 4.21
            box(slide, x, 1.45, 3.95, 5.65, fill=WHITE, line=BLACK)
            box(slide, x, 1.45, 3.95, 0.62, fill=NAVY)
            text(slide, x + 0.24, 1.6, 3.6, 0.4, t, size=15, font=MONO, color=CREAM, bold=True)
            paras = [[("✓  ", {"font": MONO, "bold": True, "color": AH}), (b, {})] for b in items]
            text(slide, x + 0.3, 2.35, 3.4, 4.6, paras, size=14.5, color=MUTE, leading=1.3)

    elif layout == "result":
        box(slide, 0.5, 1.45, 5.3, 2.6, fill=WHITE, line=BLACK)
        box(slide, 0.5, 1.45, 5.3, 0.55, fill=NAVY)
        text(slide, 0.72, 1.58, 3.6, 0.32, "DISTINCT PROVIDERS", size=13.5, font=MONO,
             color=CREAM, bold=True)
        text(slide, 3.6, 1.61, 2.0, 0.3, "STATEWIDE", size=11, font=MONO,
             color=BAND_META, bold=True, align=PP_ALIGN.RIGHT)
        logo(slide, "AH", 0.8, 2.32, h=0.4)
        text(slide, 2.75, 2.14, 2.75, 0.68, fmt(TOT_AH), size=40, font=MONO, color=AH,
             bold=True, align=PP_ALIGN.RIGHT)
        box(slide, 0.75, 3.06, 4.8, 0.022, fill=BLACK)
        logo(slide, "OH", 0.8, 3.28, h=0.45)
        text(slide, 2.75, 3.16, 2.75, 0.68, fmt(TOT_OH), size=40, font=MONO, color=OH,
             bold=True, align=PP_ALIGN.RIGHT)
        # county mini-table
        box(slide, 0.5, 4.3, 5.3, 2.9, fill=WHITE, line=BLACK)
        kicker(slide, 0.75, 4.46, 4.8, "clinicians by county", color=NAVY, size=12.5)
        rows = spec["counties"]
        ty = 4.85
        text(slide, 3.25, ty - 0.04, 1.15, 0.27, "AH", size=12, font=MONO, color=AH,
             bold=True, align=PP_ALIGN.RIGHT)
        text(slide, 4.42, ty - 0.04, 1.15, 0.27, "OH", size=12, font=MONO, color=OH,
             bold=True, align=PP_ALIGN.RIGHT)
        for name, v in rows:
            ty += 0.28
            lead_ah = v["ah"] >= v["oh"]
            text(slide, 0.75, ty, 2.55, 0.27, name, size=12.5, font=MONO, color=MUTE)
            text(slide, 3.25, ty, 1.15, 0.27, fmt(v["ah"]), size=12.5, font=MONO,
                 color=AH, bold=lead_ah, align=PP_ALIGN.RIGHT)
            text(slide, 4.42, ty, 1.15, 0.27, fmt(v["oh"]), size=12.5, font=MONO,
                 color=OH, bold=not lead_ah, align=PP_ALIGN.RIGHT)
        if spec["shot"] and Path(spec["shot"]).exists():
            box(slide, 5.92, 2.05, 6.98, 3.96, fill=WHITE, line=BLACK)
            slide.shapes.add_picture(spec["shot"], Inches(6.0), Inches(2.13),
                                     width=Inches(6.82))
        text(slide, 0.5, 7.24, 12.3, 0.3,
             f"{N_SPECS} specialties  ·  {N_ZIPS} ZIP codes  ·  {N_COUNTIES} counties  ·  every count drills to named providers",
             size=12, font=MONO, color=FAINT)

    elif layout == "h2h":
        box(slide, 0.5, 1.55, 5.3, 3.85, fill=WHITE, line=BLACK)
        kicker(slide, 0.75, 1.72, 4.8, "specialty head-to-head (statewide)", color=NAVY, size=12)
        ty = 2.12
        text(slide, 3.25, ty - 0.04, 1.15, 0.27, "AH", size=12, font=MONO, color=AH,
             bold=True, align=PP_ALIGN.RIGHT)
        text(slide, 4.42, ty - 0.04, 1.15, 0.27, "OH", size=12, font=MONO, color=OH,
             bold=True, align=PP_ALIGN.RIGHT)
        for s in spec["rows"]:
            ty += 0.4
            nm = s["name"] if len(s["name"]) <= 22 else s["name"][:21] + "…"
            lead_ah = s["ah"] >= s["oh"]
            text(slide, 0.75, ty, 2.55, 0.3, nm, size=12.5, font=MONO, color=MUTE)
            text(slide, 3.25, ty, 1.15, 0.3, fmt(s["ah"]), size=13, font=MONO,
                 color=AH, bold=lead_ah, align=PP_ALIGN.RIGHT)
            text(slide, 4.42, ty, 1.15, 0.3, fmt(s["oh"]), size=13, font=MONO,
                 color=OH, bold=not lead_ah, align=PP_ALIGN.RIGHT)
        bullets_block(slide, 0.75, 5.75, 5.0, spec["bullets"], size=14.5)
        if spec["shot"] and Path(spec["shot"]).exists():
            box(slide, 5.92, 2.05, 6.98, 3.96, fill=WHITE, line=BLACK)
            slide.shapes.add_picture(spec["shot"], Inches(6.0), Inches(2.13),
                                     width=Inches(6.82))

prs.save(DESK / "Provider-Map-Deck.pptx")
print(f"wrote {DESK / 'Provider-Map-Deck.pptx'}")

# =============================================================================
# HTML twin (pixel-true rehearsal deck: real webfonts, arrow-key navigation)
# =============================================================================
import base64
import html as html_mod

fonts_css = (ROOT / "data" / "fonts.css").read_text()
b64img = lambda p: base64.b64encode(Path(p).read_bytes()).decode()
ah64, oh64 = b64img(AH_LOGO), b64img(OH_LOGO)
county64 = b64img(COUNTY_SHOT) if COUNTY_SHOT and Path(COUNTY_SHOT).exists() else ""
zip64 = b64img(ZIP_SHOT) if ZIP_SHOT and Path(ZIP_SHOT).exists() else ""
esc = html_mod.escape

def h_bullets(items, cls="bul"):
    return "".join(f'<li>{esc(b)}</li>' for b in items)

sections = []
for spec in SLIDES:
    L = spec["layout"]
    if L == "title":
        body = f"""
      <div class="titleband"></div>
      <div class="tmark">{esc(spec['mark'][0])}<span class="blue">{esc(spec['mark'][1])}</span></div>
      <div class="tlogos"><img class="lg" src="data:image/png;base64,{ah64}"><span class="x">×</span><img class="lg lgoh" src="data:image/png;base64,{oh64}"></div>
      <div class="tsub">{esc(spec['subtitle'])}</div>
      <div class="tdate">{esc(spec['date'])}</div>"""
        sections.append(f'<section class="slide title">{body}</section>')
        continue
    bandh = f'<div class="band"><h2>{esc(spec["band"])}</h2><span>{esc(spec.get("meta",""))}</span></div>'
    if L == "stats+bullets":
        stats = ""
        for sys, num, cap, sub in spec["stats"]:
            l64 = ah64 if sys == "AH" else oh64
            c = "ah" if sys == "AH" else "oh"
            stats += f'''<div class="stat"><img src="data:image/png;base64,{l64}" class="slogo {c}"><div class="snum {c}">{esc(num)}</div><div class="scap">{esc(cap.upper())}</div><div class="ssub">{esc(sub)}</div></div>'''
        body = f'<div class="cols"><div class="statcol">{stats}</div><ul class="bul">{h_bullets(spec["bullets"])}</ul></div>'
    elif L == "bigstat+bullets":
        body = f'''<div class="cols"><div class="bigstat"><div class="bignum">{esc(spec["big"][0])}</div><div class="bigcap">{esc(spec["big"][1].upper())}</div></div><ul class="bul">{h_bullets(spec["bullets"])}</ul></div><div class="verdict">{esc(spec["verdict"])}</div>'''
    elif L == "cards" or L == "cards+proof":
        lead = f'<div class="lead">{esc(spec["lead"])}</div>' if spec.get("lead") else ""
        cards = "".join(
            f'<div class="card {"a" if i%2==0 else "o"}"><div class="ck">{esc(t)}</div><div class="cb">{esc(b)}</div></div>'
            for i, (t, b) in enumerate(spec["cards"]))
        proof = f'<div class="proof">{esc(spec["proof"])}</div>' if spec.get("proof") else ""
        body = f'{lead}<div class="grid2">{cards}</div>{proof}'
    elif L == "sequence":
        steps = "".join(
            f'<div class="step"><div class="sh">{i+1}&nbsp;&nbsp;{esc(t)}</div><div class="sb">{esc(b)}</div></div>'
            for i, (t, b) in enumerate(spec["steps"]))
        body = f'<div class="grid3">{steps}</div><div class="verdict">{esc(spec["verdict"])}</div>'
    elif L == "twopanel":
        panes = ""
        for i, (t, items) in enumerate([spec["left"], spec["right"]]):
            panes += f'<div class="pane"><div class="sh {"blue" if i else ""}">{esc(t)}</div><ul class="bul small">{h_bullets(items)}</ul></div>'
        body = f'<div class="grid2t">{panes}</div>'
    elif L == "threecol":
        cols = ""
        for t, items in spec["cols"]:
            lis = "".join(f'<li><span class="tick">✓</span>{esc(b)}</li>' for b in items)
            cols += f'<div class="pane"><div class="sh">{esc(t)}</div><ul class="checks">{lis}</ul></div>'
        body = f'<div class="grid3t">{cols}</div>'
    elif L == "result":
        rows = "".join(
            f'<tr><td>{esc(n)}</td><td class="ah{" b" if v["ah"]>=v["oh"] else ""}">{fmt(v["ah"])}</td><td class="oh{" b" if v["oh"]>v["ah"] else ""}">{fmt(v["oh"])}</td></tr>'
            for n, v in spec["counties"])
        shot = f'<div class="shot"><img src="data:image/jpeg;base64,{county64}"></div>' if county64 else ""
        body = f'''<div class="rescols"><div><div class="hcard"><div class="hband"><span>DISTINCT PROVIDERS</span><span class="m">STATEWIDE</span></div>
          <div class="hrow"><img src="data:image/png;base64,{ah64}" class="hlogo"><span class="hnum ah">{fmt(TOT_AH)}</span></div><hr>
          <div class="hrow"><img src="data:image/png;base64,{oh64}" class="hlogo oh"><span class="hnum oh">{fmt(TOT_OH)}</span></div></div>
          <div class="tbl"><div class="ck">CLINICIANS BY COUNTY</div><table><tr><th></th><th class="ah">AH</th><th class="oh">OH</th></tr>{rows}</table></div></div>{shot}</div>
          <div class="foot">{N_SPECS} specialties · {N_ZIPS} ZIP codes · {N_COUNTIES} counties · every count drills to named providers</div>'''
    elif L == "h2h":
        rows = "".join(
            f'<tr><td>{esc(s["name"])}</td><td class="ah{" b" if s["ah"]>=s["oh"] else ""}">{fmt(s["ah"])}</td><td class="oh{" b" if s["oh"]>s["ah"] else ""}">{fmt(s["oh"])}</td></tr>'
            for s in spec["rows"])
        shot = f'<div class="shot"><img src="data:image/jpeg;base64,{zip64}"></div>' if zip64 else ""
        body = f'''<div class="rescols"><div><div class="tbl"><div class="ck">SPECIALTY HEAD-TO-HEAD (STATEWIDE)</div>
          <table><tr><th></th><th class="ah">AH</th><th class="oh">OH</th></tr>{rows}</table></div>
          <ul class="bul small">{h_bullets(spec["bullets"])}</ul></div>{shot}</div>'''
    else:
        body = ""
    sections.append(f'<section class="slide">{bandh}<div class="content">{body}</div></section>')

html_out = f"""<!doctype html><html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Provider Map — Story Deck</title>
<style>
{fonts_css}
:root{{--navy:#{NAVY};--cream:#{CREAM};--cream90:#{CREAM90};--mute:#{MUTE};--faint:#{FAINT};--ah:#{AH};--oh:#{OH};
--mono:"JetBrains Mono",ui-monospace,Menlo,monospace;--disp:"Inter","Helvetica Neue",Helvetica,Arial,sans-serif}}
html,body{{margin:0;background:#222;font-family:var(--disp);color:var(--navy)}}
.slide{{width:1280px;height:720px;background:var(--cream);margin:24px auto;position:relative;overflow:hidden;
box-sizing:border-box;padding:38px 48px;display:none;box-shadow:0 4px 30px rgba(0,0,0,.5);flex-direction:column}}
.slide.on{{display:flex}}
.content{{flex:1;display:flex;flex-direction:column;justify-content:safe center;gap:22px;min-height:0;padding-top:22px}}
.band{{background:var(--navy);color:var(--cream);display:flex;justify-content:space-between;align-items:center;
padding:15px 24px;font-family:var(--mono);margin-bottom:0;flex:none}}
.band h2{{margin:0;font-size:24px;letter-spacing:.08em;font-weight:600}}
.band span{{color:#{BAND_META};font-size:15px;font-weight:600;letter-spacing:.03em}}
.cols{{display:grid;grid-template-columns:5fr 6fr;gap:40px;align-items:start}}
.stat{{background:#fff;border:2px solid #000;padding:26px 28px;margin-bottom:26px;position:relative}}
.slogo{{height:42px}}.slogo.oh{{height:48px}}
.snum{{position:absolute;top:14px;right:26px;font-family:var(--mono);font-size:76px;font-weight:800}}
.snum.ah,.ah{{color:var(--ah)}}.snum.oh,.oh{{color:var(--oh)}}
.scap{{font-family:var(--mono);font-size:16px;font-weight:700;letter-spacing:.06em;color:var(--faint);margin-top:18px}}
.ssub{{font-size:18.5px;color:var(--mute);margin-top:7px}}
ul.bul{{list-style:none;padding:0;margin:0}}
ul.bul li{{font-size:22.5px;color:var(--mute);line-height:1.5;margin-bottom:22px;padding-left:32px;position:relative}}
ul.bul li:before{{content:"—";position:absolute;left:0;color:var(--faint);font-family:var(--mono);font-weight:700}}
ul.bul.small li{{font-size:19.5px}}
.bigstat{{background:#fff;border:2px solid #000;padding:44px 30px;text-align:center;align-self:center;min-width:480px}}
.bignum{{font-family:var(--mono);font-size:170px;font-weight:800;color:var(--oh);line-height:1}}
.bigcap{{font-family:var(--mono);font-size:19px;font-weight:700;letter-spacing:.05em;margin-top:22px}}
.verdict{{background:var(--cream90);border:2px solid #000;
padding:20px 26px;font-family:var(--mono);font-weight:700;font-size:21px}}
.lead{{font-size:26px;font-weight:700;margin-bottom:0}}
.grid2{{display:grid;grid-template-columns:1fr 1fr;gap:24px}}
.card{{background:#fff;border:2px solid #000;padding:26px 28px 26px 36px;position:relative;min-height:150px;display:flex;flex-direction:column;justify-content:center}}
.card:before{{content:"";position:absolute;left:0;top:0;bottom:0;width:8px}}
.card.a:before{{background:var(--ah)}}.card.o:before{{background:var(--oh)}}
.ck{{font-family:var(--mono);font-size:18px;font-weight:700;letter-spacing:.06em;margin-bottom:12px}}
.cb{{font-size:19px;color:var(--mute);line-height:1.45}}
.proof{{background:var(--navy);color:var(--cream);
padding:22px 26px;font-family:var(--mono);font-weight:700;font-size:19.5px;line-height:1.45}}
.grid3{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:24px;align-items:stretch}}
.step{{display:flex;flex-direction:column;min-height:340px}}
.step,.pane{{background:#fff;border:2px solid #000}}
.sh{{background:var(--navy);color:var(--cream);font-family:var(--mono);font-weight:700;font-size:19px;
letter-spacing:.05em;padding:15px 22px}}
.sh.blue{{background:var(--ah)}}
.sb{{padding:26px 24px;font-size:19px;color:var(--mute);line-height:1.5;flex:1}}
.grid2t{{display:grid;grid-template-columns:1fr 1fr;gap:28px;align-items:stretch}}
.grid2t .pane{{min-height:440px}}
.grid2t ul.bul{{padding:22px 22px}}
.grid3t{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:18px}}
ul.checks{{list-style:none;margin:0;padding:24px 24px}}
ul.checks li{{font-size:17.5px;color:var(--mute);line-height:1.45;margin-bottom:16px;padding-left:30px;position:relative}}
.tick{{position:absolute;left:0;color:var(--ah);font-family:var(--mono);font-weight:700}}
.rescols{{display:grid;grid-template-columns:5fr 8fr;gap:22px;align-items:center}}
.hcard{{background:#fff;border:2px solid #000;margin-bottom:14px}}
.hband{{background:var(--navy);color:var(--cream);font-family:var(--mono);font-weight:700;font-size:15px;
display:flex;justify-content:space-between;padding:11px 18px}}
.hband .m{{color:#{BAND_META};font-size:11px}}
.hrow{{display:flex;justify-content:space-between;align-items:center;padding:13px 20px}}
.hlogo{{height:36px}}.hlogo.oh{{height:42px}}
.hnum{{font-family:var(--mono);font-size:46px;font-weight:800}}
.hcard hr{{border:none;border-top:2px solid #000;margin:0 16px}}
.tbl{{background:#fff;border:2px solid #000;padding:14px 20px}}
.tbl table{{width:100%;border-collapse:collapse;font-family:var(--mono);font-size:15px}}
.tbl th{{text-align:right;font-size:14px;padding-bottom:8px}}
.tbl td{{padding:4px 0;color:var(--mute)}}
.tbl td.ah,.tbl td.oh{{text-align:right;padding-left:16px}}
.tbl td.b{{font-weight:800}}
.shot{{background:#fff;border:2px solid #000;padding:6px;align-self:center}}
.shot img{{width:100%;display:block}}
.foot{{font-family:var(--mono);font-size:14.5px;color:var(--faint)}}
.slide.title{{padding:0}}
.titleband{{background:var(--navy);height:170px}}
.tmark{{width:640px;margin:74px auto 0;background:#fff;border:3px solid #000;text-align:center;
font-family:var(--mono);font-size:52px;font-weight:800;letter-spacing:.12em;padding:24px 0;color:var(--navy)}}
.tmark .blue{{color:var(--ah)}}
.tlogos{{display:flex;justify-content:center;align-items:center;gap:26px;margin-top:44px}}
.lg{{height:54px}}.lgoh{{height:60px}}
.x{{font-family:var(--mono);color:var(--faint);font-size:30px}}
.tsub{{text-align:center;font-size:24px;margin-top:46px}}
.tdate{{text-align:center;font-family:var(--mono);font-size:14.5px;color:var(--faint);margin-top:64px}}
.pager{{position:fixed;right:18px;bottom:14px;color:#aaa;font-family:var(--mono);font-size:12px}}
</style></head><body>
{''.join(sections)}
<div class="pager"><span id="pg"></span> · ←/→</div>
<script>
const slides=[...document.querySelectorAll('.slide')];let i=0;
const show=n=>{{i=Math.max(0,Math.min(slides.length-1,n));slides.forEach((s,j)=>s.classList.toggle('on',j===i));
document.getElementById('pg').textContent=(i+1)+' / '+slides.length;}};
addEventListener('keydown',e=>{{if(e.key==='ArrowRight'||e.key===' ')show(i+1);if(e.key==='ArrowLeft')show(i-1);}});
addEventListener('click',e=>show(i+1));show(0);
</script></body></html>"""

(DESK / "Provider-Map-Deck.html").write_text(html_out)
print(f"wrote {DESK / 'Provider-Map-Deck.html'} ({len(html_out)//1024} KB)")
print(f"slides: {len(SLIDES)} | totals AH {TOT_AH} OH {TOT_OH} | h2h rows {len(head2head)}")
